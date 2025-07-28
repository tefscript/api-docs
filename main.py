from fastapi import FastAPI, HTTPException, status, File, UploadFile
from pptx import Presentation
import docx
import base64
from pydantic import BaseModel #validação e parsing de dados
from io import BytesIO
import zipfile #lida com PPTX como arquivo ZIP
import openpyxl #arquivos excel embutidos no pptx

app = FastAPI()

#modelo para dados Base64
class FileBase64(BaseModel):
    filename: str
    data: str  #base64 do arquivo
    mimeType: str

#função auxiliar para extrair texto de DOCX
def extract_text_from_docx(file_stream: BytesIO) -> str:
    try:
        doc = docx.Document(file_stream)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip() != ""]
        return "\n".join(paragraphs)
    except Exception as e:
        raise ValueError(f"Erro DOCX: {e}")

#função auxiliar para extrair texto de PPTX
def extract_text_from_pptx(file_stream: BytesIO) -> str:
    all_extracted_text = []

    try:
        file_stream.seek(0) #reseta o ponteiro do stream antes de carregar o arquivo
        prs = Presentation(file_stream)

        for slide in prs.slides:
            texts_on_slide = []
            for shape in slide.shapes:
                #extrair o texto normal
                if hasattr(shape, "text"):
                    texts_on_slide.append(shape.text)

                #extrair texto de tabelas no powerpoint
                if shape.has_table:
                    for row in shape.table.rows:
                        row_texts = []
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    row_texts.append(paragraph.text.strip())
                        if row_texts:
                            texts_on_slide.append(" | ".join(row_texts)) #formato simples para tabela

            if texts_on_slide:
                all_extracted_text.append("\n".join(texts_on_slide))

        #extrair e processar objetos OLE (planilhas Excel embutidas no pptx)
        #resetar o ponteiro do stream novamente, pois 'prs = Presentation(file_stream)' pode ter movido.
        file_stream.seek(0)
        
        #pptx é um arquivo ZIP, abrir com zipfile
        with zipfile.ZipFile(file_stream, 'r') as zf:
            for member in zf.namelist():
                # Objetos embutidos (como planilhas Excel) geralmente estão em 'ppt/embeddings/'
                if member.startswith('ppt/embeddings/'):
                    if member.endswith('.xlsx') or member.endswith('.xls'):
                        try:
                            with zf.open(member) as embedded_file: #abre o arquivo embutido
                                excel_data = embedded_file.read()
                                #carrega a planilha com openpyxl a partir dos bytes
                                wb = openpyxl.load_workbook(BytesIO(excel_data))
                                for sheet_name in wb.sheetnames:
                                    sheet = wb[sheet_name]
                                    for row in sheet.iter_rows():
                                        row_values = [cell.value for cell in row if cell.value is not None]
                                        if row_values:
                                            all_extracted_text.append(" | ".join(map(str, row_values)))
                        except Exception as ex:
                            print(f"Erro ao processar planilha embutida '{member}': {ex}")
                            all_extracted_text.append(f"[[Erro: Não foi possível extrair conteúdo da planilha embutida '{member}']]")

    except Exception as e:
        raise ValueError(f"Erro ao processar arquivo PPTX: {e}")

    return "\n".join(all_extracted_text)

#endpoint arquivos Base64 (DOCX e PPTX)
@app.post("/extract-text-base64")
async def extract_text_base64(file: FileBase64):
    try:
        binary_data = base64.b64decode(file.data)
        file_stream = BytesIO(binary_data)

        extracted_text = ""
        if file.mimeType == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            extracted_text = extract_text_from_docx(file_stream)
        elif file.mimeType == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            extracted_text = extract_text_from_pptx(file_stream)
        else:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Tipo de MIME não suportado para extração de texto: {file.mimeType}"
            )

        return {"extracted_text": extracted_text, "file_name": file.filename, "mime_type": file.mimeType}

    except ValueError as ve:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"Falha arquivo '{file.filename}'. {ve}"
        )
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Erro interno do servidor ao processar o arquivo '{file.filename}'. Detalhes: {e}"
        )

#endpoint para arquivos binários (mantido para compatibilidade, mas o N8N usará o base64)
@app.post("/extract-pptx-binary")
async def extract_pptx_binary(file: UploadFile = File(...)):
    try:
        contents = await file.read()
        file_stream = BytesIO(contents)
        extracted_text = extract_text_from_pptx(file_stream)
        return {"extracted_text": extracted_text, "file_name": file.filename, "mime_type": file.content_type}
    except ValueError as ve:
        raise HTTPException(status_code=status.HTTP_422_UNPROCESSABLE_ENTITY, detail=f"Falha ao extrair texto do arquivo PPTX: {ve}")
    except Exception as e:
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail=f"Erro interno ao processar PPTX: {e}")

@app.post("/extract-docx-binary")
async def extract_docx_binary(file: UploadFile = File(...)):
    try:
        contents = await file.read()
        file_stream = BytesIO(contents)
        extracted_text = extract_text_from_docx(file_stream)
        return {"extracted_text": extracted_text, "file_name": file.filename, "mime_type": file.content_type}
    except ValueError as ve:
        raise HTTPException(status_code=status.HTTP_422_UNPROCESSABLE_ENTITY, detail=f"Falha ao extrair texto do arquivo DOCX: {ve}")
    except Exception as e:
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail=f"Erro interno ao processar DOCX: {e}")